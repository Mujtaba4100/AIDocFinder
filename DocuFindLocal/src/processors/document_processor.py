"""Unified document processor for many formats."""

from typing import Dict, Any, List
from pathlib import Path
import docx
import pptx
from openpyxl import load_workbook
import pandas as pd
from datetime import datetime

from src.processors.pdf_processor import PDFProcessor
from src.processors.image_processor import ImageProcessor
from src.utils.file_utils import FileUtils


class DocumentProcessor:
    def __init__(self, device: str = None):
        self.pdf_processor = PDFProcessor()
        self.image_processor = ImageProcessor(device=device)
        self.file_utils = FileUtils()
    
    def process_docx(self, docx_path: str) -> Dict[str, Any]:
        """Process DOCX file"""
        try:
            doc = docx.Document(docx_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Extract metadata
            core_properties = doc.core_properties
            metadata = {
                "title": core_properties.title or "",
                "author": core_properties.author or "",
                "subject": core_properties.subject or "",
                "keywords": core_properties.keywords or "",
                "created": core_properties.created.isoformat() if core_properties.created else "",
                "modified": core_properties.modified.isoformat() if core_properties.modified else ""
            }
            
            return {
                "file_type": "docx",
                "text": text,
                "metadata": metadata
            }
        except Exception as e:
            print(f"Error processing DOCX {docx_path}: {e}")
            return {"file_type": "docx", "text": "", "metadata": {}}
    
    def process_txt(self, txt_path: str) -> Dict[str, Any]:
        """Process TXT file"""
        try:
            with open(txt_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            return {
                "file_type": "txt",
                "text": text,
                "metadata": {}
            }
        except Exception as e:
            print(f"Error processing TXT {txt_path}: {e}")
            return {"file_type": "txt", "text": "", "metadata": {}}
    
    def process_pptx(self, pptx_path: str) -> Dict[str, Any]:
        """Process PPTX file"""
        try:
            prs = pptx.Presentation(pptx_path)
            text = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            
            return {
                "file_type": "pptx",
                "text": "\n".join(text),
                "metadata": {}
            }
        except Exception as e:
            print(f"Error processing PPTX {pptx_path}: {e}")
            return {"file_type": "pptx", "text": "", "metadata": {}}
    
    def process_xlsx(self, xlsx_path: str) -> Dict[str, Any]:
        """Process XLSX file"""
        try:
            wb = load_workbook(xlsx_path, read_only=True, data_only=True)
            text = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text.append(f"Sheet: {sheet_name}")
                
                # Read first few rows
                for row in ws.iter_rows(min_row=1, max_row=10, values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    text.append(row_text)
            
            return {
                "file_type": "xlsx",
                "text": "\n".join(text),
                "metadata": {}
            }
        except Exception as e:
            print(f"Error processing XLSX {xlsx_path}: {e}")
            return {"file_type": "xlsx", "text": "", "metadata": {}}
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """Process any supported file"""
        path = Path(file_path)
        extension = path.suffix.lower()
        
        # Get file metadata
        file_metadata = self.file_utils.get_file_metadata(file_path)
        
        # Process based on file type
        if extension == '.pdf':
            result = self.pdf_processor.process(file_path)
        elif extension in ['.doc', '.docx']:
            result = self.process_docx(file_path)
        elif extension == '.txt':
            result = self.process_txt(file_path)
        elif extension in ['.ppt', '.pptx']:
            result = self.process_pptx(file_path)
        elif extension in ['.xls', '.xlsx']:
            result = self.process_xlsx(file_path)
        elif extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif']:
            result = self.image_processor.process(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")
        
        # Add file metadata
        result["file_metadata"] = file_metadata
        
        # Generate unique ID
        result["id"] = file_metadata["file_hash"]
        
        return result
    
    def batch_process(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """Process multiple files"""
        results = []
        for file_path in file_paths:
            try:
                result = self.process_file(file_path)
                results.append(result)
                print(f"✓ Processed: {Path(file_path).name}")
            except Exception as e:
                print(f"✗ Failed to process {file_path}: {e}")
        
        return results

# Test the processor
if __name__ == "__main__":
    processor = DocumentProcessor(device="cpu")
    
    # Create test files
    test_files = []
    
    # Create a test text file
    test_txt = Path("test_document.txt")
    if not test_txt.exists():
        with open(test_txt, 'w') as f:
            f.write("Hostel Rules and Regulations\n")
            f.write("1. Curfew at 10 PM\n")
            f.write("2. No loud music after 9 PM\n")
            f.write("3. Keep rooms clean\n")
            f.write("4. Respect other students\n")
        test_files.append(str(test_txt))
    
    # Create a test PDF if pandoc is available
    test_pdf = Path("test_document.pdf")
    if not test_pdf.exists():
        try:
            import subprocess
            subprocess.run(["pandoc", "test_document.txt", "-o", "test_document.pdf"], 
                          capture_output=True)
            if test_pdf.exists():
                test_files.append(str(test_pdf))
        except:
            pass
    
    # Process files
    if test_files:
        print(f"Processing {len(test_files)} files...")
        results = processor.batch_process(test_files)
        
        for result in results:
            print(f"\nFile type: {result['file_type']}")
            if 'text' in result:
                print(f"Text preview: {result['text'][:100]}...")
            if 'description' in result:
                print(f"Description: {result['description']}")
    else:
        print("No test files created. Install pandoc for PDF creation.")
